#!/usr/bin/env python3
import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_ppt():
    # 创建一个空的演示文稿
    prs = Presentation()
    
    # 第一张幻灯片：标题页
    slide_layout = prs.slide_layouts[0]  # 标题和副标题布局
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "郭黎明 - 面试答辩"
    subtitle.text = "姓名：郭黎明\n公司：待填写\n区域：待填写\n技能：IT运维、网络管理、云计算"
    
    # 第二张幻灯片：基本资料
    slide_layout = prs.slide_layouts[1]  # 标题和内容布局
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "基本资料"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "性别：男\n民族：汉族\n年龄：27岁\n面貌：党员\n地址：无锡市滨湖区水秀新村\n电话：13673788132\n邮箱：956014740@qq.com"
    
    # 第三张幻灯片：教育背景
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "教育背景"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "时间：2018.9至2021.6\n学校：无锡城市职业技术学院\n专业：计算机网络技术\n学历：专科"
    
    # 第四张幻灯片：工作经历
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "工作经历"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "无锡飞谱电子信息技术有限公司（2021.2-2022.8）\n- 公司内部研发网络系统规划部署\n- 公司PC、服务器、交换机、防火墙、路由器等配置及运维\n- 华为云、阿里云、超算等资源的系统环境配置及运维\n- 公司设备采购；公司邮箱、网站管理；机房管理；固定资产管理\n\n嘉环科技股份有限公司（2022.9-2023.9）\n- 华东云数据中心联通软件研究院项目现场运维\n- 主机设备的硬件维修；主机夯机、宕机、失联等故障的日志分析\n- 机房安全管理；固件微码升级；防火墙工单；RAID管理配置；维修工单处理\n- 机房故障巡检；基地备品备件管理；资产核查；设备共计两千余台\n\nIT主管（2023.9至今）\n- 华东云数据中心IT主管\n- 基地主机设备维护；基地网络设备维护；动环系统维护；及客户设备故障处理"
    
    # 第五张幻灯片：专业技能
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "专业技能"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "综合布线、交换路由\n熟练掌握Linux基础命令，Windows系统\nDocker、K8s环境搭建运维、Mysql常用命令\nShell脚本、ansible集群控制\nWindows server域控管理\n私有云OpenStack和公有云的搭建运维"
    
    # 第六张幻灯片：职业证书
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "职业证书"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "RHCE红帽认证工程师\n局域网中级管理\n1+X云计算平台运维与开发初级\n云计算、大数据、综合布线、网页制作与网站开发\nSQL server、Linux高级应用、Windows server域控管理、网络安全等"
    
    # 第七张幻灯片：个人介绍
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "个人介绍"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "本人喜欢骑行，能很好地与他人沟通，具有良好的团队合作精神，工作认真踏实。\n善于与同事相处，做事踏实，能迅速地适应各种环境，并融合其中。"
    
    # 保存演示文稿
    prs.save("/home/lmguo/桌面/郭黎明-面试答辩.pptx")
    print("PPT文件已成功创建！")

if __name__ == "__main__":
    create_ppt()